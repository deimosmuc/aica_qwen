"""Build the AI Circuit Architect jury pitch deck (Track 3 — Agent Society).

Run:  .venv/Scripts/python.exe -m deck.build_deck
Out:  deck/AI_Circuit_Architect.pptx

Honest by design: the app produces an engineering *scaffold* (opens in KiCad,
structural ERC clean), not a finished/PCB-ready board. Copy avoids
"PCB-ready / manufacturable / production-ready". See the spec.
"""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import deck.theme as T

ASSETS = T.ASSETS
RGB = T.RGBColor


def _blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    T.fill_background(s)
    return s


def chip(s, x, y, text, w=Inches(2.0), h=Inches(0.42), fill=T.PANEL2,
         line=T.LINE, color=T.TEXT, size=12):
    T.panel(s, x, y, w, h, fill=fill, line=line, radius=0.5)
    T.add_text(s, x, y, w, h, text, size=size, color=color,
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def bullet(s, x, y, w, head, body, dot=T.ACCENT):
    T.badge(s, x, y + Inches(0.05), Inches(0.16), dot, "")
    T.add_text(s, x + Inches(0.32), y - Inches(0.04), w - Inches(0.32),
               Inches(0.4), head, size=15, bold=True, color=T.TEXT)
    if body:
        T.add_text(s, x + Inches(0.32), y + Inches(0.32), w - Inches(0.32),
                   Inches(0.7), body, size=13, color=T.MUTED, line_spacing=1.08)


# --------------------------------------------------------------------------- 1
def slide_title(prs):
    s = _blank(prs)
    tile = T.panel(s, Inches(5.47), Inches(1.15), Inches(2.4), Inches(2.4),
                   fill=RGB(0xFF, 0xFF, 0xFF), line=None, radius=0.18)
    tile.shadow.inherit = False
    s.shapes.add_picture(f"{ASSETS}/logo.png", Inches(5.72), Inches(1.55),
                         width=Inches(1.9))
    T.add_text(s, Inches(1.0), Inches(3.95), Inches(11.33), Inches(1.0),
               "AI Circuit Architect", size=52, bold=True, color=T.TEXT,
               align=PP_ALIGN.CENTER, font=T.FONT_H)
    T.add_text(s, Inches(1.0), Inches(5.05), Inches(11.33), Inches(0.6),
               "From a plain-English idea to a structured KiCad starting point "
               "— designed by a society of agents",
               size=20, color=T.MUTED, align=PP_ALIGN.CENTER)
    T.panel(s, Inches(4.62), Inches(6.0), Inches(4.1), Inches(0.55),
            fill=T.PANEL2, line=T.LINE, radius=0.5)
    T.add_text(s, Inches(4.62), Inches(6.0), Inches(4.1), Inches(0.55),
               "Qwen Hackathon  ·  Track: Agent Society", size=14,
               color=T.ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
               bold=True)
    return s


# --------------------------------------------------------------------------- 2
def slide_problem(prs):
    s = _blank(prs)
    T.title(s, "A schematic is slow to start — and one LLM cuts corners")

    cols = [
        ("The manual way", T.MUTED, [
            ("Blank-page problem", "Turning an idea into a structured KiCad project takes an experienced engineer."),
            ("Hours before the real work", "Sheet hierarchy, power domains, interfaces — all by hand."),
            ("Easy to forget the boring bits", "Protection, debug access, review items slip through under time pressure."),
        ]),
        ("A single LLM", T.WARN, [
            ("Fast, but shallow", "One pass gives a plausible answer — and skips the unglamorous essentials."),
            ("Misses safety / review items", "Reverse-polarity, overcurrent, reset/clock — gone in 1 of every 2 designs."),
            ("No one checks its work", "Nothing challenges the first draft or asks what's missing."),
        ]),
    ]
    cx = [Inches(0.7), Inches(7.0)]
    for (head, col, items), x in zip(cols, cx):
        T.panel(s, x, Inches(1.75), Inches(5.6), Inches(4.55), fill=T.PANEL,
                line=T.LINE)
        T.add_text(s, x + Inches(0.4), Inches(2.0), Inches(4.8), Inches(0.5),
                   head, size=18, bold=True, color=col)
        y = Inches(2.85)
        for h, b in items:
            bullet(s, x + Inches(0.4), y, Inches(4.9), h, b,
                   dot=col if col is not T.MUTED else T.LINE)
            y += Inches(1.12)
    T.add_text(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.6),
               [("Our take:  a small ", {"color": T.MUTED, "size": 15}),
                ("team of agents that argue", {"color": T.ACCENT, "size": 15, "bold": True}),
                (" → a structured, reviewed starting point you finish in KiCad.",
                 {"color": T.MUTED, "size": 15})])
    T.footer(s, 2)
    return s


# --------------------------------------------------------------------------- 3
def slide_number(prs):
    s = _blank(prs)
    T.title(s, "Two heads beat one — and we measured it")

    T.add_text(s, Inches(0.7), Inches(1.7), Inches(6.0), Inches(1.6),
               [("+3.2", {"size": 96, "bold": True, "color": T.ACCENT})],
               font=T.FONT_H)
    T.add_text(s, Inches(0.78), Inches(3.25), Inches(5.9), Inches(0.7),
               "more design concerns surfaced than a single LLM — on average, "
               "across 5 diverse hard designs (live Qwen).",
               size=15, color=T.MUTED, line_spacing=1.1)

    bar_x, bar_w_full = Inches(0.8), 5.4
    def bar(y, frac, color, label, val):
        T.hbar(s, bar_x, y, Inches(bar_w_full), Inches(0.5), T.PANEL2)
        T.hbar(s, bar_x, y, Inches(bar_w_full * frac), Inches(0.5), color)
        T.add_text(s, bar_x, y - Inches(0.34), Inches(5.4), Inches(0.3),
                   label, size=13, color=T.TEXT, bold=True)
        T.add_text(s, bar_x + Inches(bar_w_full * frac + 0.1), y + Inches(0.02),
                   Inches(1.4), Inches(0.4), val, size=14, color=color, bold=True,
                   anchor=MSO_ANCHOR.MIDDLE)
    bar(Inches(4.45), 11.6 / 12, T.ACCENT, "Multi-agent team", "11.6 / 12")
    bar(Inches(5.55), 8.4 / 12, T.MUTED, "Single agent", "8.4 / 12")
    T.add_text(s, Inches(0.8), Inches(6.35), Inches(6.0), Inches(0.5),
               "Coverage = engineering concern surfaced as work, scored by a "
               "deterministic 12-point rubric.", size=11, color=T.MUTED,
               italic=True, line_spacing=1.05)

    px, py, pw = Inches(7.2), Inches(1.7), Inches(5.45)
    T.panel(s, px, py, pw, Inches(4.75), fill=T.PANEL, line=T.LINE)
    T.add_text(s, px + Inches(0.35), py + Inches(0.28), pw - Inches(0.7),
               Inches(0.4), "Per design  ·  coverage delta (multi − single)",
               size=14, color=T.TEXT, bold=True)
    rows = [("Motor + safety", "+2", T.MUTED),
            ("Precision analog", "+2", T.MUTED),
            ("Battery IoT", "+4", T.ACCENT),
            ("Medical wearable", "+5", T.ACCENT),
            ("Industrial gateway", "+3", T.TEXT)]
    ry = py + Inches(0.95)
    for name, delta, col in rows:
        T.add_text(s, px + Inches(0.4), ry, Inches(3.6), Inches(0.5), name,
                   size=15, color=T.TEXT, anchor=MSO_ANCHOR.MIDDLE)
        T.badge(s, px + pw - Inches(1.15), ry + Inches(0.02), Inches(0.46),
                col if col != T.TEXT else T.PANEL2, delta,
                txt_color=T.BG if col is T.ACCENT else T.TEXT, size=13)
        ry += Inches(0.7)
    T.add_text(s, px + Inches(0.4), ry + Inches(0.05), pw - Inches(0.8),
               Inches(0.5),
               [("↑  the gap widens with complexity", {"color": T.ACCENT,
                 "bold": True, "size": 14})])
    T.footer(s, 3)
    return s


# --------------------------------------------------------------------------- 4
def slide_roles(prs):
    s = _blank(prs)
    T.title(s, "A team, not a prompt — six specialist agents")
    T.add_text(s, Inches(0.7), Inches(1.45), Inches(11.9), Inches(0.5),
               "Each stage has one job and its own model tier. Work hands off "
               "left to right; reviewers can send it back.", size=15,
               color=T.MUTED)

    roles = [
        ("RQ", "Requirements", "qwen-plus", T.ACCENT),
        ("AR", "Architect", "qwen-plus", T.ACCENT),
        ("DC", "Design Critic", "qwen-max", T.REVIEW),
        ("AB", "Arbitration", "qwen-max", T.WARN),
        ("PE", "PCB Engineer", "qwen-plus", T.PCB),
        ("PC", "PCB Critic", "qwen-max", T.REVIEW),
    ]
    n = len(roles)
    x0, step, cw = 0.62, 2.06, 1.78
    cy = Inches(2.5)
    for i, (ini, name, model, col) in enumerate(roles):
        x = Inches(x0 + i * step)
        T.panel(s, x, cy, Inches(cw), Inches(2.55), fill=T.PANEL, line=T.LINE)
        T.badge(s, x + Inches(cw / 2 - 0.42), cy + Inches(0.35), Inches(0.84),
                col, ini, txt_color=T.BG, size=20)
        T.add_text(s, x + Inches(0.05), cy + Inches(1.4), Inches(cw - 0.1),
                   Inches(0.6), name, size=14.5, bold=True, color=T.TEXT,
                   align=PP_ALIGN.CENTER, line_spacing=0.95)
        chip(s, x + Inches(cw / 2 - 0.62), cy + Inches(2.0), model,
             w=Inches(1.24), h=Inches(0.36),
             color=(T.REVIEW if model == "qwen-max" else T.ACCENT), size=11)
        if i < n - 1:
            T.add_text(s, x + Inches(cw - 0.02), cy + Inches(0.6),
                       Inches(0.34), Inches(0.5), "›", size=26, color=T.MUTED,
                       align=PP_ALIGN.CENTER)
    T.add_text(s, Inches(0.7), Inches(5.55), Inches(11.9), Inches(0.6),
               [("Juniors (qwen-plus) propose; ", {"color": T.MUTED, "size": 15}),
                ("senior reviewers (qwen-max) hold them to account — and demand "
                 "rework.", {"color": T.TEXT, "size": 15, "bold": True})])
    T.add_text(s, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.5),
               "A stateless orchestrator coordinates them; you can run it "
               "one-click or sign off each stage.", size=13, color=T.MUTED,
               italic=True)
    T.footer(s, 4)
    return s


# --------------------------------------------------------------------------- 5
def slide_conflict(prs):
    s = _blank(prs)
    T.title(s, "When agents disagree, the design gets better")

    # left: society-chat rework exchange, cropped to the design-side
    # negotiation (Architect's 5 blocks -> Critic flags gap -> rework -> clean)
    ct, cb = 0.06, 0.70            # crop top/bottom fractions of the capture
    img_w = Inches(6.0)
    img_h = Inches(6.0 * (1559 * (1 - ct - cb)) / 1002)
    iy = Inches(2.35)
    T.add_picture_panel(s, f"{ASSETS}/society_rework.png", Inches(0.7), iy,
                        img_w, img_h, crop=(0.0, ct, 0.0, cb), pad=0.08)
    T.add_text(s, Inches(0.77), iy + img_h + Inches(0.22), Inches(6.0),
               Inches(0.4),
               "Live Agent-Society view  ·  the Critic → Architect rework round",
               size=12, color=T.MUTED, italic=True)

    nx, nw = Inches(6.95), Inches(5.7)
    T.add_text(s, nx, Inches(1.6), nw, Inches(0.7),
               "Our system is not a pipeline. The Critic challenges the "
               "Architect, who reworks — live, in rounds.", size=16,
               color=T.MUTED, line_spacing=1.15)
    steps = [
        ("1", T.REVIEW, "Round 1 — Design Critic",
         "Flags a gap: the Architect's 5 blocks miss a Debug / SWD + status-LED block."),
        ("2", T.WARN, "↺ Rework — System Architect",
         "Reacts to the critique and revises: adds the Debug block → 6 blocks."),
        ("3", T.OK, "Round 2 — Design Critic",
         "Re-reviews the revision: no missing blocks remain. Conflict resolved."),
    ]
    sy = Inches(2.6)
    for num, col, head, body in steps:
        T.badge(s, nx, sy, Inches(0.5), col, num, txt_color=T.BG, size=16)
        T.add_text(s, nx + Inches(0.72), sy - Inches(0.04), nw - Inches(0.72),
                   Inches(0.4), head, size=15, bold=True, color=col)
        T.add_text(s, nx + Inches(0.72), sy + Inches(0.34), nw - Inches(0.72),
                   Inches(0.7), body, size=13.5, color=T.TEXT, line_spacing=1.1)
        sy += Inches(1.2)
    T.add_text(s, nx, sy + Inches(0.05), nw, Inches(0.6),
               [("This negotiation loop — not a single pass — is what raises "
                 "coverage.", {"italic": True, "color": T.ACCENT, "size": 14})],
               line_spacing=1.1)
    T.footer(s, 5)
    return s


# --------------------------------------------------------------------------- 6
def slide_efficiency(prs):
    s = _blank(prs)
    T.title(s, "Measured on a deterministic rubric — not vibes")
    T.add_text(s, Inches(0.7), Inches(1.45), Inches(11.9), Inches(0.5),
               "One 12-point engineering rubric scores both teams the same way "
               "— no LLM judging itself.", size=15, color=T.MUTED)

    # left: what the single agent skipped (real live aggregate)
    lx = Inches(0.7)
    T.panel(s, lx, Inches(2.2), Inches(7.2), Inches(4.25), fill=T.PANEL,
            line=T.LINE)
    T.add_text(s, lx + Inches(0.4), Inches(2.45), Inches(6.4), Inches(0.5),
               "What a single agent most often skipped", size=16, bold=True,
               color=T.TEXT)
    T.add_text(s, lx + Inches(0.4), Inches(2.92), Inches(6.4), Inches(0.4),
               "across the 5 live designs — caught by the Critic / Arbitration",
               size=12, color=T.MUTED, italic=True)
    items = [("Reverse-polarity protection", 5),
             ("Overcurrent / fuse protection", 4),
             ("Docs & explicit uncertainty", 3),
             ("Surge / ESD on power input", 2),
             ("Clock source", 2)]
    y = Inches(3.5)
    for label, cnt in items:
        T.add_text(s, lx + Inches(0.4), y, Inches(5.2), Inches(0.45), label,
                   size=15, color=T.TEXT, anchor=MSO_ANCHOR.MIDDLE)
        T.badge(s, lx + Inches(6.05), y + Inches(0.02), Inches(0.44), T.WARN,
                f"{cnt}×", txt_color=T.BG, size=12)
        y += Inches(0.56)

    # right: headline numbers panel
    rx = Inches(8.25)
    T.panel(s, rx, Inches(2.2), Inches(4.4), Inches(4.25), fill=T.PANEL2,
            line=T.LINE)
    T.add_text(s, rx + Inches(0.4), Inches(2.5), Inches(3.6), Inches(0.5),
               "Average coverage", size=15, bold=True, color=T.TEXT)
    T.add_text(s, rx + Inches(0.4), Inches(3.15), Inches(3.6), Inches(0.9),
               [("11.6", {"size": 40, "bold": True, "color": T.ACCENT}),
                (" / 12   multi", {"size": 16, "color": T.MUTED})])
    T.add_text(s, rx + Inches(0.4), Inches(4.05), Inches(3.6), Inches(0.9),
               [("8.4", {"size": 40, "bold": True, "color": T.MUTED}),
                (" / 12   single", {"size": 16, "color": T.MUTED})])
    T.add_text(s, rx + Inches(0.4), Inches(5.1), Inches(3.7), Inches(1.1),
               "The team's edge isn't cleverness — it's the review step that "
               "refuses to forget the unglamorous essentials.", size=13,
               color=T.TEXT, line_spacing=1.12)
    T.footer(s, 6)
    return s


# --------------------------------------------------------------------------- 7
def slide_real(prs):
    s = _blank(prs)
    T.title(s, "What you actually get — a structured KiCad start")

    # left: real generated schematic scaffold (cropped to reduce paper margin)
    T.add_picture_panel(s, f"{ASSETS}/schematic.png", Inches(0.7), Inches(1.65),
                        Inches(6.35), Inches(4.75), crop=(0.03, 0.02, 0.03, 0.03),
                        pad=0.0)
    # honest verified badge
    T.panel(s, Inches(0.9), Inches(1.85), Inches(3.2), Inches(0.5),
            fill=T.PANEL, line=T.OK, radius=0.5)
    T.add_text(s, Inches(0.9), Inches(1.85), Inches(3.2), Inches(0.5),
               [("✓  Opens in KiCad 10 · structural ERC clean",
                 {"color": T.OK, "size": 11.5, "bold": True})],
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rx, rw = Inches(7.4), Inches(5.25)
    outs = [
        ("Hierarchical KiCad scaffold", "Sheets, power-port symbols, sheet pins and an embedded block diagram — opens & passes structural checks."),
        ("PCB-readiness pack (prep)", "Net classes, constraints (.kicad_dru), candidate parts, floorplan zones, a DFT/DFM checklist."),
        ("PDF report + KiCad ZIP", "A shareable report and the full project, downloadable in one click."),
    ]
    y = Inches(1.95)
    for h, b in outs:
        T.badge(s, rx, y + Inches(0.06), Inches(0.18), T.ACCENT, "")
        T.add_text(s, rx + Inches(0.36), y - Inches(0.04), rw - Inches(0.36),
                   Inches(0.4), h, size=16.5, bold=True, color=T.TEXT)
        T.add_text(s, rx + Inches(0.36), y + Inches(0.36), rw - Inches(0.36),
                   Inches(0.8), b, size=13.5, color=T.MUTED, line_spacing=1.1)
        y += Inches(1.28)
    T.panel(s, rx, Inches(5.7), rw, Inches(0.95), fill=T.PANEL, line=T.LINE)
    T.add_text(s, rx + Inches(0.3), Inches(5.7), rw - Inches(0.6), Inches(0.95),
               [("An engineering ", {"color": T.MUTED, "size": 13}),
                ("scaffold, not a finished design", {"color": T.TEXT, "size": 13, "bold": True}),
                (" — placeholder blocks the engineer completes. No production "
                 "claims; a human stays responsible.", {"color": T.MUTED, "size": 13})],
               anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    T.footer(s, 7)
    return s


# --------------------------------------------------------------------------- 8
def slide_tech(prs):
    s = _blank(prs)
    T.title(s, "AI prepares, the human decides")

    lx, lw = Inches(0.7), Inches(5.7)
    T.add_text(s, lx, Inches(1.55), lw, Inches(0.5), "You stay in control",
               size=17, bold=True, color=T.ACCENT)
    hil = [
        ("Approve, re-run or stop — per agent", "Every stage is shown before the next one runs."),
        ("One-click or step-by-step", "Auto-run to watch the team, or sign off each stage yourself."),
        ("Audience personas", "Professional / Student / Maker re-tones every explanation."),
    ]
    y = Inches(2.2)
    for h, b in hil:
        bullet(s, lx, y, lw, h, b, dot=T.ACCENT)
        y += Inches(1.15)

    rx, rw = Inches(7.0), Inches(5.7)
    T.add_text(s, rx, Inches(1.55), rw, Inches(0.5), "Built with", size=17,
               bold=True, color=T.TEXT)
    stack = [
        ("Qwen", "qwen-plus / qwen-max, OpenAI-compatible API"),
        ("FastAPI + Alpine.js", "Python backend, live single-page UI (ELK diagrams)"),
        ("KiCad-CLI 10", "Real open / export / structural-ERC validation"),
        ("WeasyPrint · Docker", "PDF report; one-image deploy"),
    ]
    y = Inches(2.2)
    for h, b in stack:
        T.panel(s, rx, y, rw, Inches(0.92), fill=T.PANEL, line=T.LINE)
        T.add_text(s, rx + Inches(0.3), y + Inches(0.13), rw - Inches(0.6),
                   Inches(0.4), h, size=15, bold=True, color=T.TEXT)
        T.add_text(s, rx + Inches(0.3), y + Inches(0.5), rw - Inches(0.6),
                   Inches(0.36), b, size=12, color=T.MUTED)
        y += Inches(1.07)
    T.footer(s, 8)
    return s


# --------------------------------------------------------------------------- 9
def slide_close(prs):
    s = _blank(prs)
    T.title(s, "Three jury criteria, one society of agents")

    crit = [
        ("Decomposition & roles", "Six specialist agents, each with a job and a model tier."),
        ("Conflict resolution", "A real Critic → Architect rework loop, resolved in rounds."),
        ("Measured efficiency", "+3.2 coverage points over a single agent (live, 5 designs)."),
    ]
    y = Inches(1.95)
    for h, b in crit:
        T.badge(s, Inches(0.8), y, Inches(0.5), T.OK, "✓", txt_color=T.BG, size=18)
        T.add_text(s, Inches(1.5), y - Inches(0.05), Inches(6.6), Inches(0.5),
                   h, size=20, bold=True, color=T.TEXT)
        T.add_text(s, Inches(1.5), y + Inches(0.42), Inches(6.6), Inches(0.5),
                   b, size=14, color=T.MUTED)
        y += Inches(1.3)

    # right: restated number
    T.panel(s, Inches(8.7), Inches(1.95), Inches(3.95), Inches(3.6),
            fill=T.PANEL, line=T.LINE)
    T.add_text(s, Inches(8.7), Inches(2.45), Inches(3.95), Inches(1.4),
               [("+3.2", {"size": 84, "bold": True, "color": T.ACCENT})],
               align=PP_ALIGN.CENTER, font=T.FONT_H)
    T.add_text(s, Inches(8.9), Inches(3.95), Inches(3.55), Inches(1.2),
               "more engineering concerns surfaced — the case for a society of "
               "agents over a single model.", size=14, color=T.MUTED,
               align=PP_ALIGN.CENTER, line_spacing=1.12)

    T.add_text(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.5),
               [("From a plain-English idea to a structured KiCad start — "
                 "reviewed by a society of agents.",
                 {"color": T.ACCENT, "size": 15, "italic": True})])
    T.footer(s, 9)
    return s


ALL = [slide_title, slide_problem, slide_number, slide_roles, slide_conflict,
       slide_efficiency, slide_real, slide_tech, slide_close]
SAMPLE = [slide_title, slide_number, slide_conflict]


def build(slides=ALL, path="deck/AI_Circuit_Architect.pptx"):
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    for fn in slides:
        fn(prs)
    prs.save(path)
    print(f"saved {path}  ({len(slides)} slides)")


if __name__ == "__main__":
    build()
