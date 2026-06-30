"""Shared theme + layout helpers for the AI Circuit Architect pitch deck.

Palette mirrors app/static/index.html :root so the deck and product feel like
one piece. Dark throughout (premium). Motif: numbered/icon circles + rounded
panels. No accent stripes or title underlines (AI-slide tells).
"""
from __future__ import annotations
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE

# ---- palette (mirrors app/static/index.html:14-18) -------------------------
BG      = RGBColor(0x0F, 0x14, 0x19)
PANEL   = RGBColor(0x1A, 0x21, 0x2B)
PANEL2  = RGBColor(0x22, 0x2C, 0x38)
LINE    = RGBColor(0x2D, 0x3A, 0x48)
TEXT    = RGBColor(0xE6, 0xED, 0xF3)
MUTED   = RGBColor(0x8B, 0x98, 0xA5)
ACCENT  = RGBColor(0x4F, 0x9C, 0xF9)   # blue, matches logo
OK      = RGBColor(0x3F, 0xB9, 0x50)   # verified / pass
WARN    = RGBColor(0xD2, 0x99, 0x22)   # rework / conflict
REVIEW  = RGBColor(0xF7, 0x78, 0xBA)   # critic
PCB     = RGBColor(0x0D, 0x94, 0x88)   # pcb engineer teal

FONT = "Calibri"            # QA-reliable, ships with Office, ~Segoe UI
FONT_H = "Calibri"          # headers

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
ASSETS = "deck/assets"


def _no_line(shape):
    shape.line.fill.background()


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def fill_background(slide, color=BG):
    """Paint the whole slide a solid colour (full-bleed rectangle, sent back)."""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid(r, color)
    _no_line(r)
    r.shadow.inherit = False
    # send to back
    sp = r._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return r


def add_text(slide, x, y, w, h, runs, size=16, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT,
             line_spacing=1.0, italic=False):
    """Add a textbox. `runs` is a string or a list of (text, overrides) tuples.

    overrides is a dict that may set color/bold/size/italic/font per run.
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, o in runs:
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = o.get("font", font)
        f.size = Pt(o.get("size", size))
        f.bold = o.get("bold", bold)
        f.italic = o.get("italic", italic)
        f.color.rgb = o.get("color", color)
    return tb


def panel(slide, x, y, w, h, fill=PANEL, line=LINE, radius=0.08):
    """Rounded rectangle card."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _solid(shp, fill)
    if line is None:
        _no_line(shp)
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    return shp


def badge(slide, x, y, d, fill, label, txt_color=TEXT, size=14, bold=True):
    """The motif: a filled circle with centred short text (number / initials)."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    _solid(c, fill)
    _no_line(c)
    c.shadow.inherit = False
    tf = c.text_frame
    tf.word_wrap = False
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = txt_color
    return c


def hbar(slide, x, y, w, h, color, radius=0.5):
    """Horizontal value bar (used by the +3.2 chart)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _solid(shp, color)
    _no_line(shp)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    return shp


def title(slide, text, x=Inches(0.7), y=Inches(0.55), w=Inches(11.9),
          size=34, color=TEXT):
    """Slide title — weight + whitespace, no underline."""
    return add_text(slide, x, y, w, Inches(1.0), text, size=size, bold=True,
                    color=color, font=FONT_H)


def footer(slide, page_no, total=9, logo=False):
    """Muted footer: wordmark bottom-left + page number bottom-right."""
    if logo:
        try:
            slide.shapes.add_picture(f"{ASSETS}/logo.png", Inches(0.6),
                                     Inches(7.02), height=Inches(0.3))
        except Exception:
            pass
    add_text(slide, Inches(0.7), Inches(7.04), Inches(5), Inches(0.3),
             "AI Circuit Architect", size=10, color=MUTED)
    add_text(slide, Inches(11.6), Inches(7.04), Inches(1.2), Inches(0.3),
             f"{page_no} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def add_picture_panel(slide, img, x, y, w, h, crop=None, pad=0.0):
    """Place an image inside a rounded panel; crop=(l,t,r,b) fractions."""
    panel(slide, x - Inches(pad), y - Inches(pad),
          w + Inches(2 * pad), h + Inches(2 * pad), fill=PANEL2, line=LINE)
    pic = slide.shapes.add_picture(img, x, y, w, h)
    if crop:
        pic.crop_left, pic.crop_top, pic.crop_right, pic.crop_bottom = crop
    return pic
